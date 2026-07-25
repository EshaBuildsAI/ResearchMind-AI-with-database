"""
app/services/document_processor.py
Extracts raw text from uploaded documents. Only module that touches
file-parsing libraries — ported as-is from V3, per the roadmap.
"""

import io

import fitz  # PyMuPDF
import docx
import pptx
import pandas as pd

from app.utils.text import clean_text, get_file_extension


def extract_text_from_pdf(file_bytes: bytes) -> str:
    text_parts = []
    with fitz.open(stream=file_bytes, filetype="pdf") as pdf_doc:
        for page in pdf_doc:
            text_parts.append(page.get_text())
    return clean_text("\n".join(text_parts))


def extract_pdf_pages(file_bytes: bytes) -> list:
    """Returns [(page_number, text), ...], 1-indexed — used for real page-level citations."""
    pages = []
    with fitz.open(stream=file_bytes, filetype="pdf") as pdf_doc:
        for i, page in enumerate(pdf_doc, start=1):
            page_text = clean_text(page.get_text())
            if page_text:
                pages.append((i, page_text))
    return pages


def extract_text_from_docx(file_bytes: bytes) -> str:
    document = docx.Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    paragraphs.append(cell.text)
    return clean_text("\n".join(paragraphs))


def extract_text_from_pptx(file_bytes: bytes) -> str:
    presentation = pptx.Presentation(io.BytesIO(file_bytes))
    text_parts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line.strip():
                        text_parts.append(line)
    return clean_text("\n".join(text_parts))


def extract_text_from_xlsx(file_bytes: bytes) -> str:
    sheets = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
    text_parts = []
    for sheet_name, df in sheets.items():
        text_parts.append(f"Sheet: {sheet_name}")
        text_parts.append(df.to_string(index=False))
    return clean_text("\n".join(text_parts))


def extract_text_from_txt(file_bytes: bytes) -> str:
    return clean_text(file_bytes.decode("utf-8", errors="ignore"))


_EXTRACTORS = {
    "pdf": extract_text_from_pdf,
    "docx": extract_text_from_docx,
    "pptx": extract_text_from_pptx,
    "xlsx": extract_text_from_xlsx,
    "txt": extract_text_from_txt,
}


def process_document(filename: str, file_bytes: bytes) -> dict:
    """Dispatch to the correct extractor. Returns text + metadata + (PDF only) pages."""
    extension = get_file_extension(filename)
    if extension not in _EXTRACTORS:
        raise ValueError(f"Unsupported file type: .{extension}")

    text = _EXTRACTORS[extension](file_bytes)
    if not text:
        raise ValueError(f"No readable text found in '{filename}'.")

    pages = extract_pdf_pages(file_bytes) if extension == "pdf" else []

    return {
        "filename": filename,
        "extension": extension,
        "text": text,
        "char_count": len(text),
        "word_count": len(text.split()),
        "pages": pages,
    }
