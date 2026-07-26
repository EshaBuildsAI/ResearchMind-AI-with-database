"""
app/services/pptx_generator.py
Turns the LLM-generated slide outline (title + bullets per slide) into an
actual, designed .pptx file — python-pptx (free, already a dependency for
document text extraction). No paid design API; GPT-4o-mini still writes
the content, this module only handles visual layout.

Design rules followed (avoiding the generic "AI slide deck" look):
- Brand colors only (teal/coral — matches the app's own theme), never a
  default blue.
- No accent stripes/bars/underlines under titles — full background color
  blocks instead, or a circular badge for the slide number.
- Real font-size hierarchy (36-44pt titles vs 18-20pt body).
- Bullets render as colored dot + text, not plain default PowerPoint
  bullets, with generous spacing.
- Cover and closing slides get the full dark brand background; content
  slides stay white for readability.
"""

import os
import uuid

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from app.core.config import settings
from app.utils.text import ensure_dir

TEAL = RGBColor(0x14, 0xB8, 0xA6)
TEAL_DARK = RGBColor(0x0F, 0x4C, 0x4C)
CORAL = RGBColor(0xFF, 0x6F, 0x5E)
VOID = RGBColor(0x07, 0x0F, 0x0F)
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_MUTED = RGBColor(0x55, 0x60, 0x5E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _blank_slide(prs, bg_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # fully blank layout
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send background to back
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return slide


def _add_text(slide, left, top, width, height, text, size, color, bold=False,
              align=PP_ALIGN.LEFT, font="Inter", anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def _add_bullet_dot(slide, left, top, color):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.14), Inches(0.14))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    dot.shadow.inherit = False


def _add_number_badge(slide, number, left, top):
    size = Inches(0.55)
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    badge.fill.solid()
    badge.fill.fore_color.rgb = TEAL
    badge.line.fill.background()
    badge.shadow.inherit = False
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = VOID


def _cover_slide(prs, title, subtitle):
    slide = _blank_slide(prs, VOID)
    # subtle brand glow circles (not stripes) for a designed feel
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2.5), Inches(6), Inches(6))
    glow.fill.solid()
    glow.fill.fore_color.rgb = TEAL_DARK
    glow.fill.transparency = 0
    glow.line.fill.background()
    glow.shadow.inherit = False

    _add_text(slide, Inches(1), Inches(2.9), Inches(11.3), Inches(1.6), title, 44, WHITE, bold=True)
    _add_text(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.8), subtitle, 20, TEAL, bold=False)
    return slide


def _content_slide(prs, index, title, bullets):
    slide = _blank_slide(prs, WHITE)
    _add_number_badge(slide, index, Inches(0.7), Inches(0.65))
    _add_text(slide, Inches(1.5), Inches(0.6), Inches(11), Inches(0.9), title, 30, INK, bold=True)

    # bullets, sized down a bit if there are many, to avoid overflow
    n = max(len(bullets), 1)
    body_size = 20 if n <= 4 else (18 if n <= 6 else 16)
    gap = Inches(0.85) if n <= 4 else (Inches(0.65) if n <= 6 else Inches(0.5))
    y = Inches(2.0)
    for i, bullet in enumerate(bullets[:8]):  # hard cap to avoid a slide nobody can read
        text = bullet if len(bullet) <= 140 else bullet[:137].rstrip() + "..."
        _add_bullet_dot(slide, Inches(1.55), y + Inches(0.12), CORAL if i % 2 else TEAL)
        _add_text(slide, Inches(1.9), y, Inches(10.6), gap, text, body_size, INK)
        y += gap
    return slide


def _closing_slide(prs, text="Thank you"):
    slide = _blank_slide(prs, VOID)
    _add_text(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(1.2), text, 40, WHITE, bold=True, align=PP_ALIGN.CENTER)
    return slide


def build_presentation(title: str, slides: list) -> str:
    """slides: [{"title": str, "bullets": [str, ...]}, ...] (from
    llm_service.generate_presentation_outline). Returns the absolute path
    to the generated .pptx file on disk."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _cover_slide(prs, title, "Generated by ResearchMind AI")
    for i, slide_data in enumerate(slides, start=1):
        _content_slide(prs, i, slide_data.get("title", f"Slide {i}"), slide_data.get("bullets", []))
    _closing_slide(prs)

    ensure_dir(settings.EXPORT_DIR)
    filename = f"{uuid.uuid4()}.pptx"
    path = os.path.join(settings.EXPORT_DIR, filename)
    prs.save(path)
    return path